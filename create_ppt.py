from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

bg_color = RGBColor(30, 41, 59)      # Slate 800 (Elegant Dark)
title_color = RGBColor(56, 189, 248) # Sky Blue Accent
text_color = RGBColor(248, 250, 252) # White text

def apply_styling(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color

# --- SLIDE 1: MAIN TITLE ---
slide_layout = prs.slide_layouts[0] 
slide = prs.slides.add_slide(slide_layout)
apply_styling(slide)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Ambica Store"
title.text_frame.paragraphs[0].font.color.rgb = title_color
title.text_frame.paragraphs[0].font.size = Pt(64)

subtitle.text = (
    "A Premium Musical Instruments E-Commerce Platform\n\n"
    "DEVELOPED BY:\n"
    "Karina [Surname] - [Roll No.]\n"
    "Raj [Surname] - [Roll No.]\n"
    "Yug [Surname] - [Roll No.]"
)
for p in subtitle.text_frame.paragraphs:
    p.font.color.rgb = text_color
    p.font.size = Pt(22)  # Forced small size to prevent vertical overflow
    p.alignment = PP_ALIGN.CENTER

img_path = "static/images/tabla_hero_image.png"
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(9.5), Inches(2), width=Inches(3.2))

# --- SLIDE 2: ACADEMIC DETAILS ---
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
apply_styling(slide)
title = slide.shapes.title
body = slide.placeholders[1]

title.text = "Academic Details"
title.text_frame.paragraphs[0].font.color.rgb = title_color
title.text_frame.paragraphs[0].font.size = Pt(40)

tf = body.text_frame
tf.text = (
    "Program: B.Sc. / M.Sc. (CA & IT)\n"
    "Semester: 6\n\n"
    "Institution: K. S. School of Business Management & IT\n"
    "University: Gujarat University\n\n"
    "Internship Organization: [Company Name Here]\n"
    "Guide / Mentor Name: [Mentor Name Here]\n"
    "Group ID: [Insert Group ID]"
)
for p in tf.paragraphs:
    p.font.color.rgb = text_color
    p.font.size = Pt(24) # Forced size prevents vertical overflow

def add_clean_slide(prs, title_text, bullets, image_paths=None):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    apply_styling(slide)
    
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.paragraphs[0].font.size = Pt(38)
    
    tf = body.text_frame
    tf.word_wrap = True
    
    has_image = bool(image_paths)
    if has_image:
        body.width = Inches(8.5) # Protect from images overlapping text
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            
        p.font.color.rgb = text_color
        if bullet.startswith("  -"):
            p.level = 1
            p.text = bullet.strip(" -")
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(226, 232, 240)
        else:
            p.level = 0
            p.text = bullet
            p.font.size = Pt(26)

    if has_image:
        if isinstance(image_paths, str):
            image_paths = [image_paths]
            
        for idx, img_path in enumerate(image_paths):
            if os.path.exists(img_path):
                if len(image_paths) == 1:
                    slide.shapes.add_picture(img_path, Inches(9.4), Inches(2.2), width=Inches(3.3))
                else:
                    top_pos = 1.8 + (idx * 1.6)
                    slide.shapes.add_picture(img_path, Inches(10.2), Inches(top_pos), width=Inches(1.4), height=Inches(1.4))

# Content
add_clean_slide(prs, "1. Introduction to Internship", [
    "Internship Overview:",
    "  - Explored comprehensive full-stack web development.",
    "  - Hands-on experience building a modern e-commerce application.",
    "Domain:",
    "  - Software Development (Full Stack Web Architecture)",
    "Training Duration:",
    "  - 120 Hours of development."
])

add_clean_slide(prs, "2. Project Overview (Ambica Store)", [
    "Project Philosophy:",
    "  - Premium storefront for handcrafted musical instruments.",
    "Core Objectives:",
    "  - Centralize inventory securely.",
    "  - Algorithmic pricing for 'Low' to 'Premium' tiers.",
    "  - Interactive Amazon-like shopping UI."
], "static/images/Violin.png")

add_clean_slide(prs, "3. Target Audience & Business Value", [
    "Target Users (Customers):",
    "  - Professional Musicians",
    "  - Music Enthusiasts and Collecters",
    "  - Music Schools and Students",
    "Store Staff & Administrators:",
    "  - Easy-to-use digital dashboard allowing complete command over product availability."
])

add_clean_slide(prs, "4. Key Features Delivered", [
    "Secure Client Authentication System.",
    "Advanced Product Catalog functionality.",
    "Dynamic Image Swapping for color variations.",
    "Auto-calculating Secure Cart details.",
    "Complete Admin Dashboard for controls."
], "static/images/Traditional instrument.png")

add_clean_slide(prs, "5. Special Innovation: Dynamic Variants", [
    "The Problem:",
    "  - Customers couldn't visualize instrument colors accurately.",
    "The Solution:",
    "  - Removed CSS & integrated JS Source Swapping.",
    "  - Directly pulls dedicated Black/Brown/Cream images.",
    "  - Increases UI trust deeply."
], ["static/images/Black wood.jpeg", "static/images/Brown Wood.jpeg", "static/images/Cream wood.jpeg"])

add_clean_slide(prs, "6. Scope & Limitations", [
    "Current Project Scope:",
    "  - Complete localized product purchasing cycle.",
    "  - Digital warranty and returns tracking setup.",
    "Limitations:",
    "  - Payment mechanisms are restricted to manual local QR processing.",
    "  - Third-party Logistics APIs (Delivery Tracking) are not integrated."
])

add_clean_slide(prs, "7. Development Methodology", [
    "Model: Iterative Agile Methodology",
    "Phase 1: Conceptual Architecture & UI Prototyping.",
    "Phase 2: Database blueprinting & SQLite initial bridging.",
    "Phase 3: Sprints covering modular Python routing, Frontend DOM updates.",
    "Phase 4: Transition to Databases (MySQL/PostgreSQL)."
])

add_clean_slide(prs, "8. Frontend Technology Stack", [
    "Core Standards:",
    "  - HTML5 & Customized Vanilla CSS3",
    "Design Paradigm:",
    "  - Dark-mode 'Glassmorphism' (Frosted glass visual effects).",
    "  - Dynamic hover animations enhancing the premium feel.",
    "JavaScript Functionality:",
    "  - Real-time DOM manipulation for instantaneous Cart & Display updates."
])

add_clean_slide(prs, "9. Backend Development Stack", [
    "Core Language & Framework:",
    "  - Python running on Flask Web Framework.",
    "Architecture Logic:",
    "  - Modular Flask 'Blueprints' mapped into patterns.",
    "  - Robust security checkpoints.",
    "Security Mechanisms:",
    "  - Werkzeug Secure Password Hashing."
])

add_clean_slide(prs, "10. Database Architecture", [
    "Engine Setup:",
    "  - Local Development over MySQL.",
    "Entity Framework (SQLAlchemy):",
    "  - Product Tables (distinct tiers & visual variants)",
    "  - User Tables (Customer vs Admin logic).",
    "  - Logging functionality via Cart/Order levels."
], "static/images/Bansuri.png")

add_clean_slide(prs, "11. Conclusion", [
    "The Ambica Store platform successfully digitizes organic instrument trading via secure, scalable code.",
    "Through Agile structuring, the application grew into a modular, feature-rich E-Commerce suite."
])

print("Saving PPT...")
prs.save("c:/Users/Raj/Desktop/Ambica_Store_Colorful.pptx")
print("Saved PPT!")
